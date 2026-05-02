#!/usr/bin/env python3
"""Generate EdBrio Curriculum System onboarding PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_DARK = RGBColor(0x5B, 0x21, 0xB6)
PURPLE_LIGHT = RGBColor(0xED, 0xE8, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x38, 0x6B, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=BLACK, line_spacing=1.5, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Meiryo"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_section_header(slide, number, title):
    add_bg(slide, WHITE)
    # Purple accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PURPLE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Meiryo"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_textbox(slide, Inches(1.7), Inches(0.45), Inches(10), Inches(0.8),
                title, font_size=32, color=BLACK, bold=True)
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE_LIGHT
    line.line.fill.background()


# ════════════════════════════════════════════
# SLIDE 0: Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, PURPLE)

# Subtle decorative circles
for cx, cy, sz in [(11.5, 0.5, 2.5), (1.0, 5.5, 3.0), (12.0, 6.0, 1.5)]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    c.fill.solid()
    c.fill.fore_color.rgb = PURPLE_DARK
    c.line.fill.background()
    c.fill.fore_color.brightness = 0.1

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "EdBrio", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.0),
            "カリキュラムシステム 開発ガイド", font_size=36, color=WHITE, bold=False,
            alignment=PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.05))
div.fill.solid()
div.fill.fore_color.rgb = WHITE
div.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
            "新メンバー オンボーディング資料", font_size=20, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
            "2026年5月", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 1: Agenda
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "アジェンダ", font_size=36, color=BLACK, bold=True)

agenda_items = [
    ("1", "全体概要", "EdBrioとは何か"),
    ("2", "カリキュラムシステムの位置づけ", "他機能との関係性"),
    ("3", "開発環境", "技術スタック・アーキテクチャ"),
    ("4", "実装済み機能", "データモデル・コンポーネント構成"),
    ("5", "改善点", "今後実装すべき機能"),
    ("6", "今後のスケジュール", "開発マイルストーン"),
    ("7", "報酬", "報酬体系について"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.6 + i * 0.75)
    add_rounded_rect(slide, Inches(0.8), y, Inches(0.55), Inches(0.55),
                     PURPLE, num, font_size=18, font_color=WHITE)
    add_textbox(slide, Inches(1.6), y, Inches(4), Inches(0.55),
                title, font_size=20, color=BLACK, bold=True)
    add_textbox(slide, Inches(5.5), y + Inches(0.05), Inches(5), Inches(0.5),
                desc, font_size=16, color=GRAY)


# ════════════════════════════════════════════
# SLIDE 2: 全体概要 (1/2) - EdBrioとは
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "全体概要 — EdBrioとは")

add_multi_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.5), [
    "EdBrioは、教育者と学習者をつなぐ教育プラットフォームです。",
    "講師が生徒の学習を効率的に管理し、保護者と円滑にコミュニケーションできる環境を提供します。",
], font_size=20, color=BLACK, line_spacing=1.8)

# Target segments
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5), Inches(0.5),
            "ターゲットセグメント", font_size=22, color=PURPLE, bold=True)

# Box 1: 家庭教師
add_rounded_rect(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.5),
                 PURPLE_LIGHT, "", font_size=14)
add_textbox(slide, Inches(1.2), Inches(4.4), Inches(4.8), Inches(0.5),
            "🏠  家庭教師向け", font_size=22, color=PURPLE_DARK, bold=True)
add_multi_text(slide, Inches(1.2), Inches(5.0), Inches(4.8), Inches(1.5), [
    "・個人で活動する英語講師",
    "・生徒管理・スケジュール管理",
    "・保護者へのレポート送信",
    "・オンライン決済（Stripe）",
], font_size=16, color=BLACK, line_spacing=1.6)

# Box 2: 個別指導教室
add_rounded_rect(slide, Inches(6.8), Inches(4.2), Inches(5.5), Inches(2.5),
                 PURPLE_LIGHT, "", font_size=14)
add_textbox(slide, Inches(7.2), Inches(4.4), Inches(4.8), Inches(0.5),
            "🏫  個別指導教室向け", font_size=22, color=PURPLE_DARK, bold=True)
add_multi_text(slide, Inches(7.2), Inches(5.0), Inches(4.8), Inches(1.5), [
    "・複数講師の管理",
    "・教室単位での生徒管理",
    "・カリキュラムの標準化",
    "・今後展開予定",
], font_size=16, color=BLACK, line_spacing=1.6)


# ════════════════════════════════════════════
# SLIDE 3: 全体概要 (2/2) - 主要機能一覧
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "全体概要 — 主要機能一覧")

features = [
    ("📅 予約管理", "生徒がカレンダーから\n授業を予約", BLUE),
    ("💬 チャット", "講師↔保護者の\nリアルタイムメッセージ", GREEN),
    ("📝 レポート", "AIを活用した\n授業レポート自動生成", ORANGE),
    ("🎫 チケット", "授業回数券の\n購入・管理", PURPLE),
    ("📚 カリキュラム", "生徒ごとの学習計画\n目標・進捗管理", RED),
    ("💳 決済", "Stripe Connectによる\nオンライン決済", RGBColor(0x64, 0x74, 0x8B)),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.8 + row * 2.7)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(3.7), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    box.line.width = Pt(1)
    box.shadow.inherit = False

    # Color accent bar at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x + Inches(0.05), y + Inches(0.05),
                                     Inches(3.6), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
                title, font_size=20, color=color, bold=True)
    add_multi_text(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(1.2),
                   desc.split("\n"), font_size=15, color=GRAY, line_spacing=1.5)


# ════════════════════════════════════════════
# SLIDE 4: カリキュラムシステムの位置づけ
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "カリキュラムシステムの位置づけ")

add_multi_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.0), [
    "カリキュラムシステムは他の機能とは独立して動作します。",
    "講師が生徒ごとに学習計画を作成・管理する機能です。",
], font_size=20, color=BLACK, line_spacing=1.8)

# System diagram
# Center box
add_rounded_rect(slide, Inches(5.2), Inches(3.2), Inches(3.0), Inches(1.2),
                 PURPLE, "カリキュラム\nシステム", font_size=18, font_color=WHITE)

# Surrounding boxes (independent)
surrounding = [
    (Inches(1.0), Inches(3.0), "予約管理"),
    (Inches(1.0), Inches(4.8), "チャット"),
    (Inches(9.5), Inches(3.0), "レポート"),
    (Inches(9.5), Inches(4.8), "チケット / 決済"),
]
for sx, sy, label in surrounding:
    add_rounded_rect(slide, sx, sy, Inches(2.5), Inches(0.9),
                     GRAY_LIGHT, label, font_size=16, font_color=GRAY, bold=False)

# Key points
add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
            "ポイント", font_size=18, color=PURPLE, bold=True)
add_multi_text(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(1.0), [
    "✓ 講師のみが管理（保護者は閲覧しない）",
    "✓ Standardプラン以上で利用可能（プランゲート付き）",
    "✓ 現時点では他機能との連携なし（将来的に連携の可能性あり）",
], font_size=16, color=BLACK, line_spacing=1.6)


# ════════════════════════════════════════════
# SLIDE 5: 開発環境 (1/2) - アーキテクチャ
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "開発環境 — アーキテクチャ")

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
            "システム構成図", font_size=20, color=PURPLE, bold=True)

# GitHub box
add_rounded_rect(slide, Inches(0.8), Inches(2.8), Inches(3.2), Inches(2.0),
                 RGBColor(0x24, 0x29, 0x2E), "GitHub\nソースコード管理\nバージョン管理", font_size=16)

# Arrow GitHub -> Vercel
arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(4.2), Inches(3.5), Inches(1.0), Inches(0.5))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = PURPLE
arrow1.line.fill.background()

add_textbox(slide, Inches(4.0), Inches(3.0), Inches(1.5), Inches(0.4),
            "auto deploy", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Vercel box
add_rounded_rect(slide, Inches(5.4), Inches(2.8), Inches(3.2), Inches(2.0),
                 BLACK, "Vercel\nNext.js ホスティング\nフロントエンド + API", font_size=16)

# Arrow Vercel -> Supabase
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(8.8), Inches(3.5), Inches(1.0), Inches(0.5))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = GREEN
arrow2.line.fill.background()

add_textbox(slide, Inches(8.6), Inches(3.0), Inches(1.5), Inches(0.4),
            "API calls", font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Supabase box
add_rounded_rect(slide, Inches(10.0), Inches(2.8), Inches(3.2), Inches(2.0),
                 GREEN, "Supabase\nPostgreSQL データベース\n認証 + RLS + Storage", font_size=16)

# Developer box (bottom)
add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(3.2), Inches(1.5),
                 PURPLE_LIGHT, "", font_size=14)
add_textbox(slide, Inches(1.1), Inches(5.6), Inches(2.8), Inches(0.4),
            "👨‍💻 開発者（あなた）", font_size=18, color=PURPLE_DARK, bold=True)
add_multi_text(slide, Inches(1.1), Inches(6.1), Inches(2.8), Inches(0.8), [
    "ローカル開発環境",
    "git push → 自動デプロイ",
], font_size=14, color=BLACK, line_spacing=1.5)

# Arrow Developer -> GitHub
arrow3 = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                 Inches(2.1), Inches(5.0), Inches(0.5), Inches(0.5))
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = RGBColor(0x24, 0x29, 0x2E)
arrow3.line.fill.background()

add_textbox(slide, Inches(2.7), Inches(5.0), Inches(1.5), Inches(0.4),
            "git push", font_size=11, color=GRAY)


# ════════════════════════════════════════════
# SLIDE 6: 開発環境 (2/2) - 技術スタック
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "開発環境 — 技術スタック")

categories = [
    ("フロントエンド", [
        "Next.js 15 (App Router)",
        "TypeScript",
        "Tailwind CSS",
        "shadcn/ui コンポーネント",
    ], BLUE),
    ("バックエンド", [
        "Supabase (PostgreSQL 17)",
        "Row Level Security (RLS)",
        "Supabase Auth",
        "API Routes (Next.js)",
    ], GREEN),
    ("インフラ・ツール", [
        "Vercel (ホスティング)",
        "GitHub (ソース管理)",
        "Stripe Connect (決済)",
        "Anthropic Claude API (AI)",
    ], PURPLE),
    ("開発ツール", [
        "pnpm (パッケージマネージャ)",
        "Playwright (E2Eテスト)",
        "next-intl (14言語対応)",
        "ESLint + TypeScript",
    ], ORANGE),
]

for i, (cat_title, items, color) in enumerate(categories):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 2.7)

    # Card
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(5.8), Inches(2.3))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    box.line.width = Pt(1)

    # Color accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, y, Inches(0.1), Inches(2.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    add_textbox(slide, x + Inches(0.4), y + Inches(0.15), Inches(5), Inches(0.4),
                cat_title, font_size=20, color=color, bold=True)

    add_multi_text(slide, x + Inches(0.4), y + Inches(0.6), Inches(5), Inches(1.6),
                   [f"・{item}" for item in items], font_size=15, color=BLACK, line_spacing=1.5)


# ════════════════════════════════════════════
# SLIDE 7: 実装済み機能 (1/3) - データモデル
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "実装済み機能 — データモデル")

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
            "カリキュラム関連テーブル（全7テーブル）", font_size=20, color=PURPLE, bold=True)

# Main table: student_profiles
add_rounded_rect(slide, Inches(4.5), Inches(2.6), Inches(4.3), Inches(1.6),
                 PURPLE, "", font_size=14)
add_textbox(slide, Inches(4.7), Inches(2.7), Inches(4.0), Inches(0.4),
            "student_profiles", font_size=18, color=WHITE, bold=True)
add_multi_text(slide, Inches(4.7), Inches(3.1), Inches(4.0), Inches(1.0), [
    "name, grade, school, notes,",
    "native_language, teacher_id",
], font_size=13, color=RGBColor(0xE0, 0xD0, 0xFF), line_spacing=1.3)

# Child tables
child_tables = [
    (Inches(0.3), Inches(4.8), "student_goals\ngoal, profile_id", GREEN),
    (Inches(3.2), Inches(4.8), "student_weak_points\nweak_point, profile_id", ORANGE),
    (Inches(6.1), Inches(4.8), "student_strengths\nstrength, profile_id", BLUE),
    (Inches(9.0), Inches(4.8), "handover_notes\ncontent, profile_id", GRAY),
]

for cx, cy, label, color in child_tables:
    add_rounded_rect(slide, cx, cy, Inches(2.7), Inches(1.2), color, label, font_size=13)

# curriculum_units (separate branch)
add_rounded_rect(slide, Inches(2.5), Inches(6.4), Inches(3.5), Inches(0.9),
                 PURPLE_DARK, "curriculum_units\ntitle, description, target_skills, status, order", font_size=12)

# skill_assessments
add_rounded_rect(slide, Inches(7.0), Inches(6.4), Inches(3.5), Inches(0.9),
                 RED, "skill_assessments\nskill_name, level(1-5), assessed_at", font_size=12)

# Arrows (simple lines from profiles to children)
for cx in [Inches(1.6), Inches(4.5), Inches(7.4), Inches(10.3)]:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   cx, Inches(4.2), Inches(0.03), Inches(0.6))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY
    line.line.fill.background()

# Relationship labels
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(2), Inches(0.3),
            "1 : N", font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 8: 実装済み機能 (2/3) - コンポーネント構成
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "実装済み機能 — コンポーネント構成")

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
            "ファイル構成", font_size=20, color=PURPLE, bold=True)

# Hooks section
add_rounded_rect(slide, Inches(0.8), Inches(2.5), Inches(5.8), Inches(2.2),
                 PURPLE_LIGHT, "", font_size=14)
add_textbox(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(0.4),
            "📁 Hooks（データ取得・操作）", font_size=18, color=PURPLE_DARK, bold=True)
add_multi_text(slide, Inches(1.1), Inches(3.1), Inches(5.4), Inches(1.5), [
    "use-student-profiles.ts",
    "  → CRUD: 取得 / 作成 / 更新 / 削除",
    "use-student-curriculum.ts",
    "  → 目標 / ユニット / スキル評価の管理",
], font_size=14, color=BLACK, line_spacing=1.4)

# Pages section
add_rounded_rect(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(2.2),
                 PURPLE_LIGHT, "", font_size=14)
add_textbox(slide, Inches(7.3), Inches(2.6), Inches(5), Inches(0.4),
            "📁 Pages（画面）", font_size=18, color=PURPLE_DARK, bold=True)
add_multi_text(slide, Inches(7.3), Inches(3.1), Inches(5.4), Inches(1.5), [
    "curriculum/page.tsx",
    "  → 生徒プロフィール一覧（プランゲート付き）",
    "curriculum/[profileId]/page.tsx",
    "  → プロフィール詳細（タブ切替）",
], font_size=14, color=BLACK, line_spacing=1.4)

# Components section
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
            "📁 Components（UIパーツ）— 11ファイル", font_size=18, color=PURPLE_DARK, bold=True)

comp_items = [
    ("student-card", "生徒カード表示"),
    ("curriculum-profile", "プロフィール表示"),
    ("unit-list / unit-form", "ユニット一覧・作成"),
    ("goal-list / goal-form", "目標一覧・作成"),
    ("skill-list / skill-form", "スキル一覧・評価"),
    ("handover-note-list / form", "引継ぎノート"),
    ("plan-gate-curriculum", "プランゲート"),
]

for i, (comp, desc) in enumerate(comp_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(5.5 + row * 0.45)
    add_textbox(slide, x, y, Inches(3), Inches(0.4),
                f"• {comp}", font_size=13, color=BLACK, bold=True)
    add_textbox(slide, x + Inches(3.2), y, Inches(2.8), Inches(0.4),
                desc, font_size=13, color=GRAY)


# ════════════════════════════════════════════
# SLIDE 9: 実装済み機能 (3/3) - 画面フロー
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "実装済み機能 — 画面フロー")

# Flow diagram
steps = [
    ("ログイン", "講師アカウントで\nログイン", PURPLE),
    ("プラン確認", "Standardプラン\n以上が必要", ORANGE),
    ("生徒一覧", "プロフィール\n一覧表示", BLUE),
    ("生徒詳細", "タブで情報切替\n(プロフィール/ユニット)", GREEN),
    ("編集・管理", "目標/ユニット/\nスキル評価を管理", RED),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 2.6)
    y = Inches(2.2)

    add_rounded_rect(slide, x, y, Inches(2.2), Inches(2.8), WHITE, "", font_size=14)
    # Border
    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, y, Inches(2.2), Inches(2.8))
    border.fill.background()
    border.line.color.rgb = color
    border.line.width = Pt(2)

    # Step number
    add_rounded_rect(slide, x + Inches(0.7), y + Inches(0.2), Inches(0.7), Inches(0.7),
                     color, str(i + 1), font_size=20)

    add_textbox(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.0), Inches(0.4),
                title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.1), y + Inches(1.6), Inches(2.0), Inches(1.0),
                   desc.split("\n"), font_size=13, color=GRAY, line_spacing=1.4)

    # Arrow between steps
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + Inches(2.25), y + Inches(1.2), Inches(0.35), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        arrow.line.fill.background()

# Key details
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
            "技術的ポイント", font_size=18, color=PURPLE, bold=True)
add_multi_text(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1.2), [
    "✓ プランゲート: Standardプラン未満の場合、アップグレード案内を表示",
    "✓ リアルタイム: Supabase のリアルタイム機能でデータ更新を即座に反映",
    "✓ RLS: Row Level Security により、講師は自分の生徒データのみアクセス可能",
], font_size=15, color=BLACK, line_spacing=1.6)


# ════════════════════════════════════════════
# SLIDE 10: 改善点
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "改善点 — 今後実装すべき機能")

improvements = [
    ("弱点（Weak Points）UI", "DBテーブルは存在するがUIコンポーネントが未実装。\n生徒の弱点を登録・表示する画面を作成。", "高", RED),
    ("強み（Strengths）UI", "DBテーブルは存在するがUIコンポーネントが未実装。\n生徒の強みを登録・表示する画面を作成。", "高", RED),
    ("進捗ダッシュボード", "カリキュラム全体の進捗状況を可視化。\nスキル評価の推移グラフなど。", "中", ORANGE),
    ("ユニット並び替え", "ドラッグ&ドロップでカリキュラムユニットの\n順序を変更できるようにする。", "中", ORANGE),
    ("検索・フィルター", "プロフィール一覧で生徒名や学年で\n検索・絞り込みができるようにする。", "中", ORANGE),
    ("データエクスポート", "カリキュラムデータをCSV/PDFで\nエクスポートする機能。", "低", GREEN),
]

for i, (title, desc, priority, color) in enumerate(improvements):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 1.8)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(5.8), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    box.line.width = Pt(1)

    # Priority badge
    add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.8), Inches(0.4),
                     color, f"優先: {priority}", font_size=10)

    add_textbox(slide, x + Inches(1.1), y + Inches(0.15), Inches(4.5), Inches(0.4),
                title, font_size=17, color=BLACK, bold=True)

    lines = desc.split("\n")
    add_multi_text(slide, x + Inches(0.3), y + Inches(0.65), Inches(5.2), Inches(0.8),
                   lines, font_size=13, color=GRAY, line_spacing=1.4)


# ════════════════════════════════════════════
# SLIDE 11: 今後のスケジュール
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 6, "今後のスケジュール")

add_multi_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8), [
    "GW（ゴールデンウィーク）期間中に集中的に開発を進めます。",
], font_size=20, color=BLACK, line_spacing=1.8)

# Timeline
milestones = [
    ("Phase 1", "環境構築・理解", "開発環境のセットアップ\nコードベースの理解\nSupabase/Vercel/GitHub接続確認", PURPLE),
    ("Phase 2", "UI実装（弱点・強み）", "weak_points UIコンポーネント\nstrengths UIコンポーネント\n既存パターンに合わせた実装", BLUE),
    ("Phase 3", "機能追加", "進捗ダッシュボード\n検索・フィルター\nユニット並び替え", GREEN),
    ("Phase 4", "テスト・改善", "テスト実施\nバグ修正\nレビュー・リリース", ORANGE),
]

for i, (phase, title, desc, color) in enumerate(milestones):
    x = Inches(0.8 + i * 3.1)
    y = Inches(3.0)

    # Phase box
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(3.8), WHITE, "", font_size=14)
    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, y, Inches(2.8), Inches(3.8))
    border.fill.background()
    border.line.color.rgb = color
    border.line.width = Pt(2)

    add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.2), Inches(0.4),
                     color, phase, font_size=12)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.5), Inches(0.4),
                title, font_size=16, color=BLACK, bold=True)

    add_multi_text(slide, x + Inches(0.15), y + Inches(1.2), Inches(2.5), Inches(2.4),
                   desc.split("\n"), font_size=13, color=GRAY, line_spacing=1.5)

    # Arrow
    if i < len(milestones) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + Inches(2.85), y + Inches(1.6), Inches(0.25), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        arrow.line.fill.background()


# ════════════════════════════════════════════
# SLIDE 12: 報酬
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 7, "報酬")

add_rounded_rect(slide, Inches(3.5), Inches(2.5), Inches(6.3), Inches(3.0),
                 PURPLE_LIGHT, "", font_size=14)

add_textbox(slide, Inches(4.0), Inches(2.8), Inches(5.3), Inches(0.5),
            "報酬体系", font_size=24, color=PURPLE_DARK, bold=True)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(4.0), Inches(3.4), Inches(5.0), Inches(0.03))
div.fill.solid()
div.fill.fore_color.rgb = PURPLE
div.line.fill.background()

add_textbox(slide, Inches(4.0), Inches(3.7), Inches(5.3), Inches(0.6),
            "初月: 月額 10,000円 〜", font_size=28, color=BLACK, bold=True)

add_multi_text(slide, Inches(4.0), Inches(4.5), Inches(5.3), Inches(1.0), [
    "※ スキルの向上・貢献度に応じて見直し予定",
], font_size=16, color=GRAY, line_spacing=1.5)


# ════════════════════════════════════════════
# SLIDE 13: まとめ
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PURPLE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
            "まとめ", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

summary_items = [
    "EdBrioは教育プラットフォーム — カリキュラムはその中の重要な機能",
    "カリキュラムシステムは講師が生徒ごとに管理する独立した機能",
    "Next.js + Supabase + Vercel の構成で開発",
    "基本的なCRUD機能は実装済み — UI改善と新機能追加が必要",
    "GW期間中に集中開発",
]

for i, item in enumerate(summary_items):
    y = Inches(2.8 + i * 0.7)
    add_textbox(slide, Inches(2), y, Inches(9), Inches(0.6),
                f"✓  {item}", font_size=20, color=WHITE)

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
            "一緒に良いプロダクトを作りましょう！", font_size=24, color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/home/user/edbrio/curriculum_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
